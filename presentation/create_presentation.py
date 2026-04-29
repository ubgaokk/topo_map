#!/usr/bin/env python3
"""
Generate TopoMap Presentation (PowerPoint)

Usage: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Handle sub-bullets (marked with →)
        if point.startswith("→"):
            p.text = point[1:].strip()
            p.level = 1
            p.font.size = Pt(18)
        else:
            p.text = "• " + point
            p.level = 0
            p.font.size = Pt(20)
        
        p.space_after = Pt(8)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "TopoMap: Endpoint-Aware Hierarchical\nPoint-Lane Graph",
        "for Driving Scene Topology Reasoning\n\nOpenLane-V2 Benchmark"
    )
    
    # Slide 2: Motivation
    add_content_slide(prs, "Motivation & Problem", [
        "Topology reasoning is crucial for autonomous driving planning",
        "Lane connectivity determines path feasibility",
        "Current approaches have limitations:",
        "→ TopoNet: No explicit endpoint modeling",
        "→ TopoMLP: Strongly coupled to detection quality",
        "→ TopoLogic: Lacks point-level granularity",
        "Key insight: Endpoint instability is the primary bottleneck"
    ])
    
    # Slide 3: Core Contributions
    add_content_slide(prs, "Core Contributions", [
        "1. Endpoint Detection Module",
        "   Explicitly predicts lane start/end points for topology",
        "2. Point-Level Reasoning",
        "   Dense 32-point sampling along each lane centerline",
        "3. Hierarchical Point-Lane Graph",
        "   Two-level message passing: Point→Lane→Topology",
        "4. Endpoint-Aware Edge Weights",
        "   Geometry-based adjacency + learned attention"
    ])
    
    # Slide 4: Architecture Overview
    add_two_column_slide(prs, "Architecture Overview",
        "Input & Backbone",
        [
            "Multi-view camera images (6-7 cameras)",
            "ResNet-50 + FPN backbone",
            "BEVFormer-style view transformer",
            "Output: BEV features [256, 100, 200]"
        ],
        "Query Decoding",
        [
            "Lane queries: 200 queries",
            "TE queries: 100 queries",
            "6-layer transformer decoder",
            "Iterative refinement"
        ]
    )
    
    # Slide 5: Endpoint Detector Detail
    add_content_slide(prs, "Endpoint Detection Module", [
        "Input: Lane queries + Lane geometry (11 points)",
        "Outputs:",
        "→ Start point prediction [x, y, conf]",
        "→ End point prediction [x, y, conf]", 
        "→ Endpoint tokens for graph aggregation",
        "Key design:",
        "→ Shared feature encoder with geometry context",
        "→ Position-aware start/end heads",
        "→ Generates semantic endpoint embeddings"
    ])
    
    # Slide 6: Point-Lane Graph
    add_content_slide(prs, "Hierarchical Point-Lane Graph", [
        "Level 1: Point → Lane Aggregation",
        "→ Multi-head attention: lane attends to its own points",
        "→ Captures local geometry details",
        "Level 2: Lane → Lane Topology",
        "→ Endpoint-aware adjacency matrix",
        "→ GNN with successor/predecessor/self-loop edges",
        "Level 3: Cross-Level Fusion",
        "→ Gated combination of point and lane features"
    ])
    
    # Slide 7: Comparison with SOTA
    add_two_column_slide(prs, "Comparison with State-of-the-Art",
        "OpenLane-V2 subset_A (OLS)",
        [
            "TopoNet: 39.8",
            "TopoMLP: 44.1",
            "TopoLogic: 44.1",
            "TopoFormer: 46.3",
            "TopoPoint: 48.8",
            "TopoMap (Ours): 46.5+ (target)"
        ],
        "Key Differences",
        [
            "Hierarchical vs flat graph",
            "Explicit endpoint modeling",
            "Point-level aggregation",
            "Geometry-aware edges",
            "Modular design"
        ]
    )
    
    # Slide 8: Ablation Study
    add_content_slide(prs, "Ablation Study Design", [
        "Exp-1: TopoNet Baseline (reproduction)",
        "Exp-2: + Endpoint Detection only",
        "Exp-3: + Point-Lane Graph only",
        "Exp-4: Full model (endpoint + point-lanep)",
        "Expected findings:",
        "→ Point-level reasoning contributes more than endpoint alone",
        "→ Combined model exceeds individual improvements",
        "→ GNN depth and point count have diminishing returns"
    ])
    
    # Slide 9: Training Details
    add_content_slide(prs, "Training Configuration", [
        "Optimizer: AdamW (lr=2e-4, wd=0.01)",
        "Scheduler: CosineAnnealingLR (30 epochs)",
        "Loss weights:",
        "→ λ_detection = 1.0",
        "→ λ_topology = 1.0",
        "→ λ_endpoint = 0.5",
        "Mixed precision (AMP) support",
        "torch.compile() for PyTorch 2.0+"
    ])
    
    # Slide 10: Limitations
    add_content_slide(prs, "Limitations & Future Work", [
        "Current limitations:",
        "→ Topology coupled to detection quality",
        "→ No temporal modeling (frame-to-frame)",
        "→ Limited to front-view TE representation",
        "Future directions:",
        "→ Temporal scene graph propagation",
        "→ SD-Map prior integration",
        "→ Multi-modal fusion (LiDAR/radar)",
        "→ Cross-domain adaptation"
    ])
    
    # Slide 11: Conclusion
    add_content_slide(prs, "Conclusion", [
        "TopoMap: A novel endpoint-aware hierarchical topology reasoner",
        "Key innovations:",
        "→ Explicit endpoint detection for lane connectivity",
        "→ Point-level graph for fine-grained geometry",
        "→ Hierarchical message passing with gating",
        "Expected improvement: +6.7 OLS over TopoNet baseline",
        "Open-source implementation with comprehensive tests",
        "Thank you! Questions?"
    ])
    
    # Save
    output_path = os.path.join(os.path.dirname(__file__), "TopoMap_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()